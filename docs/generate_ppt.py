"""Generate iMinds Tutorials project presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "iMinds-Tutorials-Project-Presentation.pptx"

PRIMARY = RGBColor(0x1E, 0x40, 0xAF)  # blue-800
ACCENT = RGBColor(0x25, 0x63, 0xEB)   # blue-600
GRAY = RGBColor(0x4B, 0x55, 0x63)


def set_title(slide, text, subtitle=None):
    if slide.shapes.title:
        slide.shapes.title.text = text
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = PRIMARY
    if subtitle and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.8, top=1.6, width=8.5, height=5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 Title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(s, "iMinds Tutorials", "Online Learning Platform for CBSE Grade 9 & 10")

    # 2 Overview
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Project Overview")
    add_bullets(s, [
        "Full-stack online tuition platform (like Udemy for school students)",
        "Three user roles: Admin, Teacher, Student",
        "Web app + Mobile app + REST API",
        "Courses: Science & Maths (CBSE Grade 9 & 10)",
        "GitHub: github.com/kirannarayanak/iminds-tutorials",
    ])

    # 3 Problem & Solution
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Problem & Solution")
    add_bullets(s, [
        "Problem: Need a single system for enrollment, payments, content, quizzes, and attendance",
        "Solution: iMinds — role-based portal with secure login and course management",
        "Students browse courses, pay online, then access lessons and quizzes",
        "Teachers create modules, upload materials, and view student progress",
        "Admins manage users, payments, enrollments, and daily attendance",
    ])

    # 4 Tech Stack
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Technology Stack")
    add_bullets(s, [
        "Frontend (Web): Next.js 14, TypeScript, Tailwind CSS",
        "Mobile: Expo (React Native)",
        "Backend: Node.js, Express, TypeScript",
        "Database: PostgreSQL 14+",
        "Auth: JWT (access + refresh tokens)",
        "Storage: Supabase (PDFs, videos) — pluggable",
        "Payments: Mock gateway (ready for Stripe / Telr UAE)",
        "Deploy: Vercel (web) + Railway (API + DB)",
    ])

    # 5 Architecture
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "System Architecture")
    add_bullets(s, [
        "Browser / Mobile App  →  REST API (port 4000)  →  PostgreSQL",
        "Monorepo structure:",
        "  • /web — Next.js dashboard & public landing page",
        "  • /backend — Express API with role middleware",
        "  • /mobile — Expo student app",
        "  • /database — schema, migrations, seed data",
        "Public home page: choose Student / Teacher / Admin entry",
    ])

    # 6 Roles
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "User Roles & Access")
    add_bullets(s, [
        "Admin: Full platform control — teachers, courses, payments, database, attendance",
        "Teacher: Create/edit courses & modules, quizzes, materials (assigned courses only)",
        "Student: Self-register, browse catalog, pay & enroll, watch content, take quizzes",
        "Role-based menus and API permissions on every request",
    ])

    # 7 Student features
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Student Features")
    add_bullets(s, [
        "Public registration with email (Udemy-style)",
        "Course catalog with creator, teachers, module outline (preview before pay)",
        "Pay & enroll — payment required before accessing lessons",
        "Promotional discounts on new courses (10–20% launch offers)",
        "Course recommendations based on enrolled subjects",
        "Dashboard: courses, quizzes, payments, class schedule",
        "In-app notifications (new courses, promotions, suggestions)",
    ])

    # 8 Teacher features
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Teacher Features")
    add_bullets(s, [
        "Create new courses (auto-assigned as instructor)",
        "Add / edit / delete modules per course",
        "Upload PDFs and materials; add video URLs",
        "Build quizzes with multiple-choice questions",
        "View assigned students and course list",
        "Receive notifications when new students register",
    ])

    # 9 Admin features
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Admin Features")
    add_bullets(s, [
        "Manage teachers (create accounts, reset passwords)",
        "Manage all courses, assign teachers, edit schedules",
        "Payment management — track paid / pending / revenue summary",
        "Database tab: enrollments + payment status per student",
        "Daily class attendance — mark present / absent by course",
        "Notifications when students sign up",
    ])

    # 10 Payments
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Payments & Enrollment Flow")
    add_bullets(s, [
        "1. Student selects course → Pay & enroll",
        "2. API creates payment (pending) — mock or real gateway",
        "3. On payment success → automatic course enrollment",
        "4. Free courses enroll immediately",
        "5. Promo discounts applied from notifications (30-day window)",
        "Admin can mark payments paid manually",
    ])

    # 11 Notifications
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Notifications System")
    add_bullets(s, [
        "New student registers → notify all admins & teachers",
        "New course published → notify all students with discount offer",
        "Smart course suggestions (e.g. Science → suggest Maths)",
        "Bell icon in header with unread count",
        "Stored in PostgreSQL notifications table",
    ])

    # 12 Database
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Database Highlights")
    add_bullets(s, [
        "PostgreSQL with 20+ tables",
        "Users, roles, student/teacher profiles",
        "Courses, modules, materials, videos, quizzes",
        "Enrollments, payments, class schedules",
        "Class attendance (daily roll call)",
        "Notifications, audit logs, refresh tokens",
    ])

    # 13 Demo credentials
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Demo Accounts (After Seed)")
    add_bullets(s, [
        "Admin:   admin / admin@123#",
        "Teacher: priyas / priyas@123#",
        "Student: aryank / aryank@123#",
        "Local web: http://localhost:3000",
        "API: http://localhost:4000/api/v1",
        "Docker: docker-compose up postgres -d",
    ])

    # 14 Deployment
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Deployment & SEO")
    add_bullets(s, [
        "Code on GitHub — pushed to main branch",
        "Web → Vercel (root folder: web)",
        "API + DB → Railway",
        "SEO: sitemap.xml, robots.txt, meta tags, JSON-LD",
        "Google Search Console — submit sitemap after go-live",
        "Guides: docs/DEPLOY.md, docs/GOOGLE_SEARCH.md",
    ])

    # 15 Future
    s = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(s, "Future Enhancements")
    add_bullets(s, [
        "Live Stripe / Telr payment integration (UAE)",
        "Live video classes (Zoom/Meet integration)",
        "Email/SMS notifications",
        "Custom domain + SSL in production",
        "Parent portal for progress reports",
        "Certificate generation on course completion",
    ])

    # 16 Thank you
    s = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(s, "Thank You", "iMinds Tutorials — Questions?")

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
